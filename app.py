__import__('pysqlite3')
import sys
sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')

import os
import streamlit as st
import tempfile
from typing import List, Dict, Optional
from io import BytesIO
import pandas as pd
import uuid
from datetime import datetime
import random
import matplotlib.pyplot as plt

# Additional libraries from Code A
from PIL import Image, ImageDraw

# Install necessary packages as needed:
# !pip install crewai streamlit PyPDF2 python-docx openpyxl python-pptx requests beautifulsoup4

from crewai import Agent, Task, Crew, Process, LLM
#from crewai_tools import SerperDevTool, FileReadTool, WebsiteSearchTool
from pydantic import BaseModel, Field
# Import FirecrawlScrapeWebsiteTool instead of SerperDevTool
from crewai_tools import FileReadTool, WebsiteSearchTool, FirecrawlScrapeWebsiteTool

# Add the time module which was missing earlier
import time

# File processing imports
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

# Set your API keys (replace with your keys or load securely)
os.environ["OPENAI_API_KEY"] = "sk-proj-123456"
#os.environ["FIRECRAWL_API_KEY"] = st.secrets["firecrawl_key"]  # Replace with your 
#os.environ["GEMINI_API_KEY"] = st.secrets["gemini_key"] # Replace with your key

# Configure Gemini LLM
gemini_llm = LLM(
    model="gemini/gemini-2.0-flash-lite",
    temperature=0.7,
)

# ---------------- Helper Functions ----------------

def generate_avatar(name, color):
    """Generate a simple circular avatar with initials."""
    name = name or "Unknown"
    color = color or "#4CAF50"
    initials = "".join([n[0].upper() for n in name.split() if n])
    if not initials:
        initials = "?"
    initials = initials[:2]
    img = Image.new('RGB', (100, 100), color=color)
    d = ImageDraw.Draw(img)
    d.ellipse((5, 5, 95, 95), fill=color)
    d.text((50, 50), initials, fill="white", anchor="mm")
    return img

# ---------------- File Processor ----------------

class FileProcessor:
    """Process various file types and extract text."""
    @staticmethod
    def process_pdf(file_bytes):
        with BytesIO(file_bytes) as f:
            reader = PyPDF2.PdfReader(f)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    @staticmethod
    def process_docx(file_bytes):
        with BytesIO(file_bytes) as f:
            doc = docx.Document(f)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    @staticmethod
    def process_xlsx(file_bytes):
        with BytesIO(file_bytes) as f:
            wb = openpyxl.load_workbook(f)
            text = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    text += " | ".join([str(cell) for cell in row if cell is not None]) + "\n"
            return text

    @staticmethod
    def process_pptx(file_bytes):
        with BytesIO(file_bytes) as f:
            prs = Presentation(f)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

    @staticmethod
    def process_txt(file_bytes):
        return file_bytes.decode('utf-8')

    @staticmethod
    def process_url(url):
        try:
            scraper = FirecrawlScrapeWebsiteTool(
                api_key=os.environ.get("FIRECRAWL_API_KEY"),
                url=url,
                page_options={"onlyMainContent": True}
            )
            result = scraper.run()
            return result
        except Exception as e:
            print(f"Error scraping URL: {e}")
            # Fallback to basic request
            response = requests.get(url)
            soup = BeautifulSoup(response.content, 'html.parser')
            return soup.get_text()

    @classmethod
    def process_file(cls, file, file_type=None):
        if file_type is None:
            file_name = file.name.lower()
            if file_name.endswith('.pdf'):
                return cls.process_pdf(file.read())
            elif file_name.endswith('.docx'):
                return cls.process_docx(file.read())
            elif file_name.endswith('.xlsx') or file_name.endswith('.xls'):
                return cls.process_xlsx(file.read())
            elif file_name.endswith('.pptx') or file_name.endswith('.ppt'):
                return cls.process_pptx(file.read())
            elif file_name.endswith('.txt'):
                return cls.process_txt(file.read())
            else:
                return "Unsupported file format"
        else:
            if file_type == 'pdf':
                return cls.process_pdf(file.read())
            elif file_type == 'docx':
                return cls.process_docx(file.read())
            elif file_type in ['xlsx', 'xls']:
                return cls.process_xlsx(file.read())
            elif file_type in ['pptx', 'ppt']:
                return cls.process_pptx(file.read())
            elif file_type == 'txt':
                return cls.process_txt(file.read())
            else:
                return "Unsupported file format"

# ---------------- Persona Model ----------------

class Persona(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    name: str
    background: str
    expertise: str
    viewpoints: str
    communication_style: str = Field(default="Professional")
    assertiveness: int = Field(default=5)  # Scale 1-10
    cooperation: int = Field(default=5)    # Scale 1-10
    color: str = Field(default="#3366cc")    # For avatar generation
    is_user: bool = Field(default=False)

# ---------------- Board Discussion System ----------------

class BoardDiscussionSystem:
    """Agentic board discussion system with dynamic personas and chat-like interface."""
    
    def __init__(self):
        # Tools
        self.web_scraper = FirecrawlScrapeWebsiteTool(
            api_key=os.environ.get("FIRECRAWL_API_KEY"),
            page_options={"onlyMainContent": True}
        )
        self.file_tool = FileReadTool(
            config=dict(
                llm=dict(
                    provider="google",
                    config=dict(
                        model="gemini-2.0-flash-lite",
                    ),
                ),
                embedder=dict(
                    provider="google",
                    config=dict(
                        model="gemini-embedding-exp-03-07",
                        task_type="retrieval_document",
                    ),
                ),
            )
        )
        self.web_tool = WebsiteSearchTool()
        
        # Base agent for research tasks
        self.research_agent = self._create_research_agent()
        
        # Dynamic attributes
        self.personas: List[Persona] = []
        self.persona_agents: List[Agent] = []
        self.user_persona: Optional[Persona] = None
        self.discussion_topic: str = ""
        self.materials: List[Dict] = []
        self.research_findings: str = ""
        self.discussion_history: List[Dict] = []
        self.discussion_dynamics: Dict[str, int] = {
            'pace': 50,         # 0 (detailed) to 100 (brief)
            'creativity': 50    # 0 (conventional) to 100 (highly creative)
        }

    def _create_research_agent(self):
        return Agent(
            role="Research Analyst",
            goal="Extract and summarize key information from documents and online sources.",
            backstory="""You are a skilled research analyst who excels at analyzing complex information 
            from various sources and summarizing key insights.""",
            verbose=True,
            tools=[self.web_scraper, self.file_tool, self.web_tool],
            llm=gemini_llm
        )
    
    def create_persona_agent(self, persona: Persona):
        return Agent(
            role=f"{persona.name} - Board Member",
            goal=f"Participate in a board discussion as {persona.name} with a unique perspective.",
            backstory=f"""You are {persona.name} with a background in {persona.background}.
            Your expertise is in {persona.expertise} and you hold the following viewpoints: {persona.viewpoints}.
            Your communication style is {persona.communication_style} with assertiveness {persona.assertiveness}/10 
            and cooperation {persona.cooperation}/10.
            Discussion Dynamics - Pace: {self.discussion_dynamics['pace']}, Creativity: {self.discussion_dynamics['creativity']}.""",
            verbose=True,
            llm=gemini_llm
        )
    
    def add_persona(self, persona: Persona):
        self.personas.append(persona)
        agent = self.create_persona_agent(persona)
        self.persona_agents.append(agent)
        return agent
    
    def set_user_persona(self, persona: Persona):
        persona.is_user = True
        self.user_persona = persona
        self.add_persona(persona)
    
    def process_materials(self, files=None, urls=None):
        processed = []
        if files:
            for file in files:
                content = FileProcessor.process_file(file)
                processed.append({
                    "name": file.name,
                    "type": "file",
                    "content": content
                })
        if urls:
            for url in urls:
                content = FileProcessor.process_url(url)
                processed.append({
                    "name": url,
                    "type": "url",
                    "content": content
                })
        self.materials = processed
        return processed

    def auto_generate_personas(self, num_personas: int) -> List[Persona]:
        """Generate detailed personas with unique personalities based on materials."""
        if not self.materials:
            return []
            
        # Create sample content from materials
        sample_content = "\n\n".join([f"{mat['name']}:\n{mat['content'][:1000]}" for mat in self.materials[:3]])
        
        # Create a prompt for generating diverse personas
        prompt = f"""Based on the following materials about "{self.discussion_topic}", 
        create {num_personas} diverse board member personas with distinct personalities, backgrounds, and viewpoints.
        
        For each persona, provide:
        1. A realistic full name
        2. Detailed professional background
        3. Specific area of expertise related to {self.discussion_topic}
        4. Unique viewpoints and perspectives they would bring to the discussion
        5. Communication style (Professional, Academic, Direct, Diplomatic, or Technical)
        6. Assertiveness level (1-10)
        7. Cooperation level (1-10)
        
        Ensure the personas have contrasting personalities and perspectives to create a dynamic discussion.
        
        MATERIALS SAMPLE:
        {sample_content}
        
        Format as a JSON array with fields: name, background, expertise, viewpoints, communication_style, assertiveness, cooperation"""
        
        task = Task(
            description=prompt,
            expected_output="JSON array of detailed persona descriptions",
            agent=self.research_agent
        )
        
        crew = Crew(
            agents=[self.research_agent],
            tasks=[task],
            verbose=True,
            process=Process.sequential
        )
        
        try:
            result = crew.kickoff()
            
            # Extract JSON from the result
            import re
            import json
            
            # Try to find JSON array in the response
            json_match = re.search(r'\[\s*\{.*\}\s*\]', result.raw, re.DOTALL)
            
            if json_match:
                json_str = json_match.group(0)
                personas_data = json.loads(json_str)
            else:
                # If no JSON array found, attempt to parse the entire result
                personas_data = json.loads(result.raw)
            
            # Create personas from parsed data
            personas = []
            for data in personas_data:
                # Generate a vibrant color for the avatar
                color = "#" + ''.join([random.choice('0123456789ABCDEF') for _ in range(6)])
                
                persona = Persona(
                    name=data.get('name', f"Board Member {len(personas)+1}"),
                    background=data.get('background', "Business professional"),
                    expertise=data.get('expertise', f"Expert in {self.discussion_topic}"),
                    viewpoints=data.get('viewpoints', "Has balanced perspectives"),
                    communication_style=data.get('communication_style', "Professional"),
                    assertiveness=min(max(int(data.get('assertiveness', 5)), 1), 10),
                    cooperation=min(max(int(data.get('cooperation', 5)), 1), 10),
                    color=color
                )
                personas.append(persona)
            
            return personas
            
        except Exception as e:
            print(f"Error generating personas: {e}")
            # Fallback to basic personas with distinct traits
            return [
                Persona(
                    name=f"{random.choice(['Dr.', 'Prof.', 'Ms.', 'Mr.'])} {random.choice(['Smith', 'Johnson', 'Lee', 'Garcia', 'Chen', 'Kumar', 'Müller', 'Rodriguez'])} {random.choice(['A.', 'B.', 'C.', 'D.', 'E.'])}",
                    background=f"Professional with {random.randint(5, 30)} years of experience in {random.choice(['finance', 'technology', 'healthcare', 'education', 'consulting', 'manufacturing'])}",
                    expertise=f"Expert in {random.choice(['strategic planning', 'digital transformation', 'risk management', 'innovation', 'operations', 'market analysis'])} related to {self.discussion_topic}",
                    viewpoints=f"Believes that {random.choice(['innovation is key', 'careful analysis is essential', 'people come first', 'efficiency drives success', 'sustainability matters most', 'adaptability is crucial'])} when discussing {self.discussion_topic}",
                    communication_style=random.choice(["Professional", "Academic", "Direct", "Diplomatic", "Technical"]),
                    assertiveness=random.randint(3, 8),
                    cooperation=random.randint(3, 8),
                    color="#" + ''.join([random.choice('0123456789ABCDEF') for _ in range(6)])
                ) for i in range(num_personas)
            ]

    def research_task(self, topic, specific_questions=None):
        materials_text = "\n\n".join(
            [f"--- {mat['name']} ---\n{mat['content'][:1000]}..." for mat in self.materials]
        )
        questions = ""
        if specific_questions:
            questions = "\n".join([f"- {q}" for q in specific_questions.split("\n") if q.strip()])
            questions = f"\n\nSpecific Questions:\n{questions}"
        prompt = f"""Analyze the following materials about: {topic}.
{questions}

MATERIALS:
{materials_text}

Provide a comprehensive analysis covering key facts, perspectives, consensus, disagreements, and implications."""
        task = Task(
            description=prompt,
            expected_output="A detailed analysis report.",
            agent=self.research_agent
        )
        return task

    def create_discussion_task(self, persona_agent, topic, research_findings, other_personas, previous_contributions=None):
        others = "\n".join([f"- {p.name}: {p.background} (Expertise: {p.expertise})" for p in other_personas])
        prev = f"PREVIOUS CONTRIBUTIONS:\n{previous_contributions}" if previous_contributions else ""
        prompt = f"""Participate in a board discussion on: {topic}

RESEARCH FINDINGS:
{research_findings}

{prev}

OTHER BOARD MEMBERS:
{others}

Respond as your persona with a perspective influenced by your background and expertise.
Discussion Dynamics:
- Pace: {self.discussion_dynamics['pace']}
- Creativity: {self.discussion_dynamics['creativity']}

Keep your response concise but thorough (150-350 words)."""
        task = Task(
            description=prompt,
            expected_output="A thoughtful board discussion contribution.",
            agent=persona_agent
        )
        return task

    def handle_user_input(self, user_input, topic, research_findings, discussion_history):
        history_text = "\n\n".join([f"{entry['persona']}: {entry['content']}" for entry in discussion_history])
        if not self.user_persona:
            return "Error: User persona not set."
        # Format the user input using the research agent (or user agent)
        prompt = f"""The user ({self.user_persona.name}) has provided this input during a board discussion:
{user_input}

Format this input to match the user's persona characteristics and integrate it smoothly into the ongoing discussion.
User Persona:
- Name: {self.user_persona.name}
- Background: {self.user_persona.background}
- Expertise: {self.user_persona.expertise}
- Viewpoints: {self.user_persona.viewpoints}
- Communication Style: {self.user_persona.communication_style}

Discussion History:
{history_text}

Return only the formatted contribution."""
        task = Task(
            description=prompt,
            expected_output="Formatted user contribution.",
            agent=self.persona_agents[[p.id for p in self.personas].index(self.user_persona.id)]
        )
        crew = Crew(
            agents=[task.agent],
            tasks=[task],
            verbose=True,
            process=Process.sequential
        )
        result = crew.kickoff()
        return result.raw

    def get_ai_responses(self, topic, research_findings, discussion_history, responding_ids: List[str]):
        responses = []
        history_text = "\n\n".join([f"{entry['persona']}: {entry['content']}" for entry in discussion_history])
        for i, persona in enumerate(self.personas):
            if persona.id not in responding_ids:
                continue
            agent = self.persona_agents[i]
            prompt = f"""Continue the board discussion on: {topic}

RESEARCH FINDINGS:
{research_findings}

Discussion History:
{history_text}

Respond as {persona.name} based on your expertise and viewpoints.
Discussion Dynamics:
- Pace: {self.discussion_dynamics['pace']}
- Creativity: {self.discussion_dynamics['creativity']}

Keep your response concise but thorough (150-350 words)."""
            task = Task(
                description=prompt,
                expected_output=f"A thoughtful response from {persona.name}.",
                agent=agent
            )
            crew = Crew(
                agents=[agent],
                tasks=[task],
                verbose=True,
                process=Process.sequential
            )
            result = crew.kickoff()
            responses.append({
                "persona_id": persona.id,
                "persona_name": persona.name,
                "content": result.raw
            })
        return responses

    def run_initial_discussion(self, topic, rounds=1):
        if self.materials:
            research_task = self.research_task(topic)
            crew = Crew(
                agents=[self.research_agent],
                tasks=[research_task],
                verbose=True,
                process=Process.sequential
            )
            research_result = crew.kickoff()
            self.research_findings = research_result.raw
        else:
            self.research_findings = "No research materials provided."
        for round_num in range(1, rounds + 1):
            for i, agent in enumerate(self.persona_agents):
                if self.personas[i].is_user:
                    continue
                others = [p for j, p in enumerate(self.personas) if j != i and not p.is_user]
                prev = "\n\n".join([f"{entry['persona']}: {entry['content']}" for entry in self.discussion_history])
                task = self.create_discussion_task(agent, topic, self.research_findings, others, previous_contributions=prev)
                crew = Crew(
                    agents=[agent],
                    tasks=[task],
                    verbose=True,
                    process=Process.sequential
                )
                result = crew.kickoff()
                self.discussion_history.append({
                    "persona_id": self.personas[i].id,
                    "persona": self.personas[i].name,
                    "content": result.raw
                })
        return {
            "research": self.research_findings,
            "discussion": self.discussion_history
        }

    def extract_key_points(self):
        if len(self.discussion_history) < 3:
            return ["Discussion in progress..."]
        conversation = "\n\n".join([f"{entry['persona']}: {entry['content']}" for entry in self.discussion_history[-20:]])
        prompt = f"""Extract 3-5 key points from this board discussion on "{self.discussion_topic}".
Format each point as a concise bullet point.

Discussion transcript:
{conversation}

Key points (as bullet points):
"""
        task = Task(
            description=prompt,
            expected_output="Key points as bullet points.",
            agent=self.research_agent
        )
        crew = Crew(
            agents=[self.research_agent],
            tasks=[task],
            verbose=True,
            process=Process.sequential
        )
        result = crew.kickoff()
        points = result.raw.split("\n")
        return [p.strip().lstrip("•-*").strip() for p in points if p.strip()]

    def generate_markdown_export(self):
        md = f"# Board Discussion: {self.discussion_topic}\n\n"
        md += f"*Date: {datetime.now().strftime('%Y-%m-%d')}*\n\n"
        md += "## Participants\n\n"
        for p in self.personas:
            md += f"- **{p.name}**: {p.background}\n"
        md += "\n## Research Findings\n\n" + self.research_findings + "\n\n"
        md += "## Discussion Transcript\n\n"
        for entry in self.discussion_history:
            md += f"**{entry['persona']}**: {entry['content']}\n\n"
        key_points = self.extract_key_points()
        if key_points:
            md += "## Key Points\n\n"
            for point in key_points:
                md += f"- {point}\n"
        return md

def handle_user_message(system, message):
    """Process user message and generate AI responses sequentially."""
    # Format user message
    with st.spinner("💭 Processing your message..."):
        formatted_message = system.handle_user_input(
            message,
            system.discussion_topic,
            system.research_findings,
            system.discussion_history
        )
        
        # Add formatted user message to history
        system.discussion_history.append({
            "persona_id": system.user_persona.id,
            "persona": system.user_persona.name,
            "content": formatted_message,
            "timestamp": datetime.now()
        })
        
        st.rerun()  # Show the user message first
    
    # Determine which personas will respond
    responding_personas = random.sample(
        [p for p in system.personas if not p.is_user],
        min(random.randint(1, 3), len(system.personas) - 1)
    )
    
    # Process responses one at a time
    for persona in responding_personas:
        with st.spinner(f"💭 {persona.name} is typing..."):
            # Add a realistic typing delay based on message length
            time.sleep(st.session_state.message_delay)
            
            # Get response from this specific persona
            response = system.get_ai_responses(
                system.discussion_topic,
                system.research_findings,
                system.discussion_history,
                [persona.id]
            )
            
            if response:
                # Add response to history
                system.discussion_history.append({
                    "persona_id": response[0]["persona_id"],
                    "persona": response[0]["persona_name"],
                    "content": response[0]["content"],
                    "timestamp": datetime.now()
                })
                
                # Update the UI to show this message
                st.rerun()

# ---------------- Streamlit App Interface ----------------

def create_streamlit_app():
    st.set_page_config(page_title="🎯 Board Discussion Simulator", layout="wide")
    
    # Initialize session state variables if they don't exist
    if 'system' not in st.session_state:
        st.session_state.system = BoardDiscussionSystem()
    if 'chat_turn' not in st.session_state:
        st.session_state.chat_turn = 0
    if 'user_joined' not in st.session_state:
        st.session_state.user_joined = False
    if 'message_delay' not in st.session_state:
        st.session_state.message_delay = 3
    if 'active_tab' not in st.session_state:
        st.session_state.active_tab = 0  # Default to first tab
        
    system: BoardDiscussionSystem = st.session_state.system

    # Sidebar Configuration
    with st.sidebar:
        st.title("⚙️ Configuration")
        
        # API Keys section
        with st.expander("🔑 API Keys"):
            #openai_key = st.text_input("OpenAI API Key", value=os.getenv("OPENAI_API_KEY"), type="password")
            firecrawl_key = st.text_input("Firecrawl API Key", value=os.getenv("FIRECRAWL_API_KEY"), type="password")
            gemini_key = st.text_input("Gemini API Key", value=os.getenv("GEMINI_API_KEY"), type="password")
            if st.button("💾 Save API Keys"):
                #os.environ["OPENAI_API_KEY"] = openai_key
                os.environ["FIRECRAWL_API_KEY"] = firecrawl_key
                os.environ["GEMINI_API_KEY"] = gemini_key
                st.success("✅ API keys saved!")

        # Discussion Settings
        st.header("🎛️ Discussion Settings")
        system.discussion_dynamics['pace'] = st.slider("⚡ Discussion Pace", 0, 100, 50)
        system.discussion_dynamics['creativity'] = st.slider("🎨 Creativity Level", 0, 100, 50)
        st.session_state.message_delay = st.slider("⏱️ Message Delay (seconds)", 1, 10, 3)
        
        # Analysis and Export
        st.header("📊 Discussion Analysis")
        if system.discussion_history:
            key_points = system.extract_key_points()
            with st.expander("🎯 Key Points", expanded=True):
                for point in key_points:
                    st.write(f"• {point}")
            
            # Participation Stats
            stats = {}
            for entry in system.discussion_history:
                stats[entry['persona']] = stats.get(entry['persona'], 0) + 1
            
            if stats:
                with st.expander("📈 Participation Stats", expanded=True):
                    fig, ax = plt.subplots(figsize=(8, 4))
                    names = list(stats.keys())
                    counts = list(stats.values())
                    ax.barh(names, counts)
                    ax.set_xlabel("Messages")
                    plt.tight_layout()
                    st.pyplot(fig)

            if st.button("📥 Export Discussion", use_container_width=True):
                md = system.generate_markdown_export()
                st.download_button(
                    "📄 Download Markdown",
                    data=md,
                    file_name=f"discussion_{datetime.now().strftime('%Y%m%d')}.md",
                    mime="text/markdown",
                    use_container_width=True
                )

        # Reset Session button at bottom of sidebar
        st.markdown("---")
        if st.button("🔄 Reset Session", use_container_width=True):
            # Clear session state
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            # Initialize a new system
            st.session_state.system = BoardDiscussionSystem()
            st.session_state.chat_turn = 0
            st.session_state.user_joined = False
            st.session_state.message_delay = 3
            st.session_state.active_tab = 0
            st.success("✨ Session reset successfully!")
            st.rerun()
    
    # Main content area
    st.title("🎯 Board Discussion Simulator")
    
    # Tab navigation
    tab_titles = ["🎯 Setup", "📚 Materials", "👥 Personas", "👤 Your Persona", "💬 Discussion"]
    current_tab = st.session_state.active_tab
    
    # Display horizontal tab buttons
    tabs = st.tabs(tab_titles)
    
    # Setup Tab
    with tabs[0]:
        st.header("🎯 Discussion Setup")
        topic = st.text_area("Discussion Topic", key="setup_topic", 
                             value=system.discussion_topic or "",
                             placeholder="Enter the main topic for the board discussion...")
        specific_questions = st.text_area("❓ Specific Questions (Optional)", 
                            help="Enter one question per line",
                            key="setup_questions",
                            placeholder="Question 1\nQuestion 2\n...")

        num_initial_rounds = st.slider("🔄 Initial Discussion Rounds", 1, 3, 1)
        
        if st.button("✅ Save Setup", use_container_width=True):
            if topic:
                system.discussion_topic = topic
                st.success("✨ Discussion topic saved!")
                time.sleep(1)
                # Auto navigate to Materials tab
                st.session_state.active_tab = 1
                st.rerun()
            else:
                st.error("❌ Please enter a discussion topic")
                
    # Materials Tab
    with tabs[1]:
        st.header("📚 Research Materials")
        st.info("📝 Upload materials to enable automatic persona generation")
        
        material_type = st.radio("📎 Source Type", 
                               ["📄 Files", "🔗 URLs", "📑 Both"])
        
        uploaded_files = None
        urls = None
        
        if material_type in ["📄 Files", "📑 Both"]:
            uploaded_files = st.file_uploader(
                "Upload Files",
                accept_multiple_files=True,
                type=["pdf", "docx", "xlsx", "xls", "pptx", "ppt", "txt"],
                help="Upload research documents"
            )
            if uploaded_files:
                st.success(f"📎 {len(uploaded_files)} files uploaded")
        
        if material_type in ["🔗 URLs", "📑 Both"]:
            urls = st.text_area(
                "Enter URLs (one per line)",
                help="Add web pages for research",
                key="material_urls"
            )
            if urls:
                url_list = [u.strip() for u in urls.split("\n") if u.strip()]
                st.success(f"🔗 {len(url_list)} URLs added")
        
        if st.button("📥 Process Materials", use_container_width=True):
            url_list = [u.strip() for u in urls.split("\n") if u.strip()] if urls else None
            with st.spinner("🔄 Processing materials..."):
                system.process_materials(
                    files=uploaded_files if uploaded_files else None,
                    urls=url_list
                )
            st.success("✨ Materials processed successfully!")
            time.sleep(1)
            
            if system.materials:
                st.subheader("📑 Processed Materials")
                for mat in system.materials:
                    with st.expander(f"{'📄' if mat['type']=='file' else '🔗'} {mat['name']}"):
                       st.text_area(
                            "Preview",
                            mat['content'][:500] + "...",
                            height=100,
                            disabled=True,
                            key=f"preview_{mat['name']}"
                        )
                
                # Auto navigate to personas tab
                st.session_state.active_tab = 2
                st.rerun()
                
    # Personas Tab
    with tabs[2]:
        st.header("👥 Define Board Member Personas")
        
        # Check if materials are available
        if not system.materials:
            st.warning("⚠️ Please upload and process materials first to enable automatic persona generation")
            st.info("💡 You can still create personas manually")
        
        # Display current personas if any exist
        if system.personas:
            st.subheader("Current Personas")
            num_personas = len([p for p in system.personas if not p.is_user])
            if num_personas > 0:
                persona_cols = st.columns(min(4, num_personas))
                idx = 0
                for persona in system.personas:
                    if not persona.is_user:  # Don't show user persona here
                        with persona_cols[idx % len(persona_cols)]:
                            st.image(generate_avatar(persona.name, persona.color), width=80)
                            st.markdown(f"**{persona.name}**")
                            with st.expander("Details"):
                                st.write(f"🎓 **Background:** {persona.background}")
                                st.write(f"🔍 **Expertise:** {persona.expertise}")
                                st.write(f"💭 **Viewpoint:** {persona.viewpoints}")
                                st.write(f"🗣️ **Style:** {persona.communication_style}")
                                st.write(f"⚡ **Assertiveness:** {'▪️' * persona.assertiveness}")
                                st.write(f"🤝 **Cooperation:** {'▪️' * persona.cooperation}")
                        idx += 1
                st.markdown("---")
        
        creation_method = st.radio(
            "Choose Method",
            ["🤖 Auto Generate", "✍️ Create Manually"]
        )

        if creation_method == "🤖 Auto Generate":
            if not system.materials:
                st.error("❌ Upload and process materials first to use automatic generation")
            else:
                num_personas = st.slider("Number of Personas", 2, 6, 4)
                if st.button("🎲 Generate Personas", use_container_width=True):
                    with st.spinner("🔄 Generating personas based on materials..."):
                        generated = system.auto_generate_personas(num_personas)
                        if generated:
                            for persona in generated:
                                system.add_persona(persona)
                            st.success(f"✨ Generated {len(generated)} personas!")
                            time.sleep(1)
                            # Auto navigate to user persona tab
                            st.session_state.active_tab = 3
                            st.rerun()
                        else:
                            st.error("❌ Could not generate personas from materials")

        else:
            st.subheader("✍️ Add New Persona")
            # Use a form with submit button
            with st.form(key="manual_persona_form"):
                col1, col2 = st.columns(2)
                with col1:
                    name = st.text_input("Name")
                    background = st.text_input("Background")
                    expertise = st.text_input("Expertise")
                with col2:
                    viewpoints = st.text_area("Key Viewpoints", key="persona_viewpoints")
                    communication_style = st.selectbox(
                        "Communication Style",
                        ["Professional", "Academic", "Direct", "Diplomatic", "Technical"]
                    )
                    assertiveness = st.slider("Assertiveness", 1, 10, 5)
                    cooperation = st.slider("Cooperation", 1, 10, 5)
                
                # Form submit button
                submit_button = st.form_submit_button("➕ Add Persona", use_container_width=True)

            if submit_button:
                if all([name, background, expertise, viewpoints]):
                    new_persona = Persona(
                        name=name,
                        background=background,
                        expertise=expertise,
                        viewpoints=viewpoints,
                        communication_style=communication_style,
                        assertiveness=assertiveness,
                        cooperation=cooperation,
                        color="#" + ''.join([random.choice('0123456789ABCDEF') for _ in range(6)])
                    )
                    system.add_persona(new_persona)
                    st.success(f"✨ Added persona: {name}")
                    time.sleep(1)
                    st.rerun()
                else:
                    st.error("❌ Please fill in all fields")
                    
        # Navigation buttons
        st.markdown("---")
        if st.button("Next: Create Your Persona →", use_container_width=True):
            st.session_state.active_tab = 3
            st.rerun()

    # User Persona Tab
    with tabs[3]:
        st.header("👤 Create Your Persona")
        st.info("🎯 Create your persona to join the discussion")
        
        if system.user_persona:
            st.success(f"✨ Your persona: {system.user_persona.name}")
            col1, col2 = st.columns([1, 3])
            with col1:
                st.image(generate_avatar(system.user_persona.name, 
                                      system.user_persona.color), 
                        width=150, caption="Your Avatar")
            with col2:
                st.write(f"**Name:** {system.user_persona.name}")
                st.write(f"**Background:** {system.user_persona.background}")
                st.write(f"**Expertise:** {system.user_persona.expertise}")
                st.write(f"**Viewpoints:** {system.user_persona.viewpoints}")
                st.write(f"**Style:** {system.user_persona.communication_style}")
                st.write(f"**Assertiveness:** {'🔵' * system.user_persona.assertiveness}")
                st.write(f"**Cooperation:** {'🟢' * system.user_persona.cooperation}")
            
            st.markdown("---")
            col1, col2 = st.columns(2)
            with col1:
                if st.button("🔄 Reset Persona", use_container_width=True):
                    for i, persona in enumerate(system.personas):
                        if persona.is_user:
                            system.personas.pop(i)
                            system.persona_agents.pop(i)
                            break
                    system.user_persona = None
                    st.success("✨ User persona reset")
                    time.sleep(1)
                    st.rerun()
            with col2:
                if st.button("Next: Start Discussion →", use_container_width=True):
                    st.session_state.active_tab = 4
                    st.rerun()
        else:
            # User form with proper submit button
            with st.form(key="user_persona_form"):
                st.subheader("✍️ Create Your Profile")
                col1, col2 = st.columns(2)
                with col1:
                    user_name = st.text_input("Your Name")
                    user_background = st.text_input("Your Background")
                    user_expertise = st.text_input("Your Area of Expertise")
                with col2:
                    user_viewpoints = st.text_area("Your Key Viewpoints")
                    user_comm_style = st.selectbox(
                        "Your Communication Style",
                        ["Professional", "Academic", "Direct", "Diplomatic", "Technical"]
                    )
                    user_assertiveness = st.slider("Your Assertiveness", 1, 10, 5,
                                                 help="How strongly you express opinions")
                    user_cooperation = st.slider("Your Cooperation", 1, 10, 5,
                                               help="How collaborative you are")
                
                # Form submit button
                persona_submit = st.form_submit_button("✨ Create My Persona", use_container_width=True)
            
            if persona_submit:
                if all([user_name, user_background, user_expertise, user_viewpoints]):
                    user_persona = Persona(
                        name=user_name,
                        background=user_background,
                        expertise=user_expertise,
                        viewpoints=user_viewpoints,
                        communication_style=user_comm_style,
                        assertiveness=user_assertiveness,
                        cooperation=user_cooperation,
                        color="#" + ''.join([random.choice('0123456789ABCDEF') for _ in range(6)]),
                        is_user=True
                    )
                    system.set_user_persona(user_persona)
                    st.success(f"✨ User persona created: {user_name}")
                    time.sleep(1)
                    st.rerun()
                else:
                    st.error("❌ Please fill in all fields")
                    
    # Discussion Tab
    with tabs[4]:
        st.header("💬 Board Discussion")
        
        # Check prerequisites
        if not system.discussion_topic:
            st.warning("⚠️ Please set up a discussion topic first")
        elif not system.materials:
            st.warning("⚠️ Please upload some research materials first")
        elif len(system.personas) < 2:
            st.warning("⚠️ Please create at least 2 personas first")
        else:
            # Display current personas
            st.subheader("🎭 Current Participants")
            num_personas = len(system.personas)
            if num_personas > 0:
                persona_cols = st.columns(min(4, num_personas))
                for i, persona in enumerate(system.personas):
                    with persona_cols[i % len(persona_cols)]:
                        st.image(generate_avatar(persona.name, persona.color), width=60)
                        st.markdown(f"**{persona.name}**" + (" 👤" if persona.is_user else ""))
            
            # Discussion Area
            st.markdown("---")
            
            if not st.session_state.get('discussion_started', False):
                st.info("🎯 Press 'Start Discussion' to begin the board meeting")
                # Make this button more prominent
                if st.button("🎬 Start Discussion", key="start_discussion", use_container_width=True):
                    with st.spinner("🔄 Starting discussion..."):
                        st.session_state.discussion_started = True
                        # Initialize discussion with first round
                        results = system.run_initial_discussion(
                            system.discussion_topic,
                            rounds=1
                        )
                    st.success("✨ Discussion started!")
                    st.rerun()
            else:
                # Show research findings at the top
                with st.expander("📚 Research Summary", expanded=False):
                    st.markdown(system.research_findings)
                
                # Message display area
                st.subheader("💬 Discussion Messages")
                
                # Create chat container
                chat_container = st.container(height=400, border=True)
                with chat_container:
                    # Display each message with avatar
                    for entry in system.discussion_history:
                        is_user = system.user_persona and (system.user_persona.name == entry['persona'])
                        persona_color = next((p.color for p in system.personas if p.name == entry['persona']), "#888888")
                        
                        with st.chat_message(
                            "user" if is_user else "assistant",
                            avatar=generate_avatar(entry['persona'], persona_color)
                        ):
                            st.markdown(f"**{entry['persona']}** {'👤' if is_user else ''}")
                            st.markdown(entry['content'])
                            st.caption(f"⏰ {entry.get('timestamp', datetime.now()).strftime('%H:%M:%S')}")
                
                # Discussion control buttons
                if not st.session_state.get('discussion_concluded', False):
                    col1, col2 = st.columns(2)
                    with col1:
                        if st.button("⏩ Continue Discussion", key="continue_discussion", use_container_width=True):
                            with st.spinner("🔄 Generating next round of discussion..."):
                                # Add a new round of AI responses
                                results = system.run_initial_discussion(
                                    system.discussion_topic,
                                    rounds=1
                                )
                            st.success("✨ Discussion continued with new responses!")
                            st.rerun()
                    with col2:
                        if st.button("🏁 Conclude Discussion", key="conclude_discussion", use_container_width=True):
                            with st.spinner("🔄 Concluding discussion..."):
                                # Generate conclusion
                                conclusion_task = Task(
                                    description=f"""As a facilitator, summarize the key points of this board discussion on "{system.discussion_topic}".
                                    Highlight areas of consensus, disagreement, and next steps.
                                    Keep it concise but comprehensive, around 200-300 words.""",
                                    expected_output="A balanced conclusion summarizing the discussion",
                                    agent=system.research_agent
                                )
                                crew = Crew(
                                    agents=[system.research_agent],
                                    tasks=[conclusion_task],
                                    verbose=True,
                                    process=Process.sequential
                                )
                                conclusion = crew.kickoff()
                                
                                # Add conclusion to history
                                system.discussion_history.append({
                                    "persona_id": "facilitator",
                                    "persona": "Meeting Facilitator",
                                    "content": f"**DISCUSSION CONCLUSION**\n\n{conclusion.raw}",
                                    "timestamp": datetime.now()
                                })
                                
                                # Mark discussion as concluded
                                st.session_state.discussion_concluded = True
                            st.success("✨ Discussion successfully concluded!")
                            st.rerun()
                else:
                    # Show message that discussion is concluded
                    st.success("🏁 This discussion has been concluded. You can export the results or start a new session.")
                    if st.button("📊 View Analysis", use_container_width=True):
                        # Auto-expand the key points section in the sidebar
                        st.session_state.show_analysis = True
                        # Change tab to analysis (optional)
                        st.rerun()
                
                # User participation section
                st.markdown("---")
                if system.user_persona:
                    if not st.session_state.get('user_joined', False):
                        # Join button
                        st.info("👤 Click below to join the discussion as your persona")
                        if st.button("👋 Join Discussion", key="join_discussion", use_container_width=True):
                            st.session_state.user_joined = True
                            st.success("✨ You've joined the discussion!")
                            st.rerun()
                    else:
                        # Chat input for user messages
                        st.markdown(f"**Chat as {system.user_persona.name}**")
                        
                        # Use a form for user input
                        with st.form(key="chat_form"):
                            user_message = st.text_area("Your message:", 
                                                    key="user_message_input", 
                                                    height=100,
                                                    placeholder=f"Type your message as {system.user_persona.name}...")
                            
                            col1, col2 = st.columns([3, 1])
                            with col1:
                                send_message = st.form_submit_button("Send Message", use_container_width=True)
                            with col2:
                                leave_discussion = st.form_submit_button("👋 Leave", use_container_width=True)
                        
                        if send_message and user_message:
                            process_user_message(system, user_message)
                            st.rerun()
                            
                        if leave_discussion:
                            st.session_state.user_joined = False
                            st.info("You've left the discussion.")
                            st.rerun()
                            
                else:
                    st.warning("👤 You need to create your persona before joining the discussion")
                    if st.button("➡️ Create My Persona", key="goto_persona_tab", use_container_width=True):
                        st.session_state.active_tab = 3
                        st.rerun()

# Define a function to handle user message processing
def process_user_message(system, message):
    # Format and add user message
    formatted_message = system.handle_user_input(
        message,
        system.discussion_topic,
        system.research_findings,
        system.discussion_history
    )
    
    # Add to history
    system.discussion_history.append({
        "persona_id": system.user_persona.id,
        "persona": system.user_persona.name,
        "content": formatted_message,
        "timestamp": datetime.now()
    })
    
    # Determine responding personas
    num_responses = random.randint(1, min(3, len(system.personas) - 1))
    responding_personas = random.sample(
        [p for p in system.personas if not p.is_user],
        num_responses
    )
    
    # Generate responses one by one
    for persona in responding_personas:
        # Get response for this persona (without delay as it's handled by streamlit)
        response = system.get_ai_responses(
            system.discussion_topic,
            system.research_findings,
            system.discussion_history,
            [persona.id]
        )
        
        if response:
            # Add to history
            system.discussion_history.append({
                "persona_id": response[0]["persona_id"],
                "persona": response[0]["persona_name"],
                "content": response[0]["content"],
                "timestamp": datetime.now()
            })

if __name__ == "__main__":
    create_streamlit_app()
